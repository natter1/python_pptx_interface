"""
This module provides an easier Interface to create *.pptx presentations using the module python-pptx.
@author: Nathanael Jöhrmann
"""
import io
import os
from pathlib import Path, PurePath
from typing import Type, Optional, Iterable, Union

from pptx_tools import utils
from pptx_tools.position import PPTXPosition
from pptx_tools.table_style import PPTXTableStyle

try:
    from matplotlib.figure import Figure
    import matplotlib.pyplot as plt

    has_matplotlib = True
except ImportError as e:
    has_matplotlib = False

import pptx
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.slide import Slide, SlideLayout
from pptx.text.text import _Run
from pptx.util import Inches

from pptx_tools.font_style import PPTXFontStyle
from pptx_tools.templates import AbstractTemplate


class PPTXCreator:
    """
    This Class provides an easy interface to create a PowerPoint presentation.
        - PPTXPosion is used to position new shapes (allowing position as fraction of slide height/width)
        - use pptx templates (in combination with templates.py)
        - removes unused placeholder from added slides
    """

    # disable typechecker, because None values are not allowed for attributes, but needed to create them in __init__
    # Correct values are set when calling self._create_presentation()
    # noinspection PyTypeChecker
    def __init__(self, template: Optional[AbstractTemplate] = None):
        self.slides: list = []
        self.template: Type[AbstractTemplate] = None
        self.prs: Presentation = None
        self.title_layout: SlideLayout = None
        self.default_layout: SlideLayout = None
        self._create_presentation(template)
        self.default_position = PPTXPosition(presentation=self.prs)

    def _fraction_width_to_inch(self, fraction: float) -> Inches:
        """
        Returns a width in inches calculated as a fraction of total slide-width.
        """
        return Inches(self.prs.slide_width.inches * fraction)

    def _fraction_height_to_inch(self, fraction: float) -> Inches:
        """
        Returns a height in inches calculated as a fraction of total slide-height.
        """
        return Inches(self.prs.slide_height.inches * fraction)

    def _create_presentation(self, template=None) -> None:
        """
        Create a new presentation (using optional template).
        """
        if template:
            self._create_presentation_from_template(template)
        else:
            self.prs = pptx.Presentation()
            self.title_layout = self.prs.slide_masters[0].slide_layouts[0]
            self.default_layout = self.prs.slide_masters[0].slide_layouts[0]

    def _create_presentation_from_template(self, template: AbstractTemplate) -> None:
        """Create a new presentation using the given template."""
        self.template = template
        self.prs = template.prs
        self.title_layout = template.title_layout
        self.default_layout = template.default_layout

    def add_title_slide(self, title: str, layout: SlideLayout = None) -> Slide:
        """Add a new slide to presentation. If no layout is given, title_layout is used."""
        if not layout:
            layout = self.title_layout
        return self.add_slide(title, layout)

    def add_slide(self, title: str, layout: SlideLayout = None) -> Slide:
        """Add a new slide to presentation. If no layout is given, default_layout is used."""
        if not layout:
            layout = self.default_layout
        slide = self.prs.slides.add_slide(layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        self.remove_unpopulated_shapes(slide)
        return slide

    def add_image(self, file: Union[Path, io.BytesIO], slide: Slide,
                  position: PPTXPosition = None,
                  zoom: float = 1.0,
                  **kwargs) -> Picture:
        """
        Add an image from disk or io.BytesIO() to slide, and position it via position.
        Optional parameter zoom sets image scaling in PowerPoint. Only used if width not in kwargs (default = 1.0).
        """
        # python-pptx (v0.6.18) can not handle Path object
        if isinstance(file, PurePath):
            file = str(file)

        if not position:
            position = self.default_position
        kwargs.update(position.dict())

        pic = slide.shapes.add_picture(file, **kwargs)  # 0, 0)#, left, top)
        pic.width = round(pic.width * zoom)
        pic.height = round(pic.height * zoom)
        return pic

    def add_matplotlib_figure(self, fig: 'Figure', slide: Slide,
                              position: PPTXPosition = None,
                              zoom: float = 1.0,
                              **kwargs) -> Picture:
        """
        Add a motplotlib figure to slide and position it via position.
        Optional parameter zoom sets image scaling in PowerPoint. Only used if width not in kwargs (default = 1.0).
        """
        if not has_matplotlib:
            raise ModuleNotFoundError("Adding a matplotlib figure needs module matplotlib to be installed.")

        with io.BytesIO() as output:
            fig.savefig(output, format="png")
            # pic = slide.shapes.add_picture(output, **kwargs)  # 0, 0)#, left, top)
            pic = self.add_image(output, slide, position, zoom, **kwargs)  # 0, 0)#, left, top)
        return pic

    def add_latex_formula(self, formula: str, slide: Slide, position: PPTXPosition = None, dpi: int = 150,
                          font_size: int = 18, color: str = "black", alpha: float = 0.0, **kwargs) -> Picture:
        """
        Add the given latex-like math-formula as an image to the presentation using matplotlib.
        """
        if not has_matplotlib:
            raise ModuleNotFoundError("Adding a latex-like formula needs module matplotlib to be installed.")

        figure: plt.Figure = plt.figure(figsize=(20, 20), dpi=dpi)
        figure.suptitle(fr"${formula}$", fontsize=font_size, color=color, **kwargs)
        tight_bbox = figure.get_tightbbox(figure.canvas.get_renderer())  # tight_layout = True
        # there seems to be no way to add the tight_bbox to existing figure ^^ -> create figure again
        figure: plt.Figure = plt.figure(figsize=tight_bbox.size, dpi=dpi)
        figure.suptitle(fr"${formula}$", fontsize=font_size, color=color, **kwargs)
        figure.patch.set_alpha(alpha)
        return self.add_matplotlib_figure(figure, slide, position)

    def add_text_box(self, slide, text: str, position: PPTXPosition = None, font: PPTXFontStyle = None) -> Shape:
        """
        Add a text box with given text using given position and paragraph.
        Uses self.default_position if no position is given.
        """
        width = height = Inches(1)  # no auto-resizing of shape -> has to be done inside PowerPoint
        if position is None:
            position = self.default_position
        result = slide.shapes.add_textbox(**position.dict(), width=width, height=height)
        result.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        result.text_frame.text = text  # first paragraph
        if font:
            font.write_shape(result)
        return result

    def _get_rows_cols(self, table_data: Iterable[Iterable[any]]):
        """Used to get number of rows and cols from table data."""
        rows = sum(1 for e in table_data)

        cols = 0
        for row in table_data:
            length = sum(1 for e in row)
            cols = max(cols, length)

        return rows, cols

    def add_table(self, slide: Slide, table_data: Iterable[Iterable[any]], position: PPTXPosition = None,
                  table_style: PPTXTableStyle = None, auto_merge: bool = False) -> Shape:
        """
        Add a table shape with given table_data at position using table_style.
        table_data: outer iter -> rows, inner iter cols
        auto_merge: use 'merge_left' and 'merge_up' as entry to mark merging cells (not implemented jet)
        """
        rows, cols = self._get_rows_cols(table_data)
        if position is None:
            position = self.default_position
        left, top = position.tuple()
        result = slide.shapes.add_table(rows, cols, left, top, width=Inches(cols), height=Inches(0.5 * rows))

        table = result.table
        for ir, row in enumerate(table_data):
            for ic, entry in enumerate(row):
                table.cell(ir, ic).text = f"{entry}"

        if table_style:
            table_style.write_shape(result)

        if auto_merge:
            pass  # todo: merge cells; replace text for merged cells with ""

        return result

    def move_slide(self, slide: Slide, new_index: int):
        """Move the given slide to position new_index."""
        _sldIdLst = self.prs.slides._sldIdLst

        to_move = None
        for entry in _sldIdLst.sldId_lst:
            if entry.id == slide.slide_id:
                to_move = entry

        if to_move is not None:
            _sldIdLst.insert(new_index, to_move)

    @staticmethod
    def remove_unpopulated_shapes(slide: Slide):
        """
        Removes empty placeholders (e.g. due to layout) from slide.
        Further testing needed.
        """
        for index in reversed(range(len(slide.shapes))):
            shape = slide.shapes[index]
            # if shape.is_placeholder and shape.text_frame.text == "":
            if shape.has_text_frame and shape.text_frame.text == "":
                shape.element.getparent().remove(shape.element)

    @staticmethod
    def create_hyperlink(run: _Run, shape: Shape, to_slide: Slide):  # text hyperlink not implemented in pptx-python
        """Make the given run a hyperlink to to_slide."""
        shape.click_action.target_slide = to_slide
        run.hyperlink.address = shape.click_action.hyperlink.address
        run.hyperlink._hlinkClick.action = shape.click_action.hyperlink._hlink.action
        run.hyperlink._hlinkClick.rId = shape.click_action.hyperlink._hlink.rId
        shape.click_action.target_slide = None

    def add_content_slide(self, slide_index=1):
        """Add a content slide with hyperlinks to all other slides and puts it to position slide_index."""
        content_entries = [
            (slide.shapes.title.text, slide) for slide in self.prs.slides
        ]


        result = self.add_slide("Content")
        content_text_box = self.add_text_box(result, "", PPTXPosition(0.1, 0.2))
        for text, slide in content_entries[1:]:
            paragraph = content_text_box.text_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = text
            self.create_hyperlink(run, content_text_box, slide)

        self.move_slide(result, slide_index)

        return result

    def save(self, filename: Union[str, "LocalPath"], create_pdf: bool = False, overwrite=False):
        """
        Save presentation under the given filename.
        """
        if os.path.isfile(filename) and not overwrite:
            print(f"File {filename} already exists. Set overwrite=True, if you want to overwrite file.")
        else:
            self.prs.save(filename)

        if create_pdf:
            filename = str(filename)  # enables to work with LocalPath-variable (which is not subscriptable)
            self.save_as_pdf(filename[:-4] + "pdf", overwrite)

    def save_as_pdf(self, filename: str, overwrite=False) -> bool:
        """
        Save the presentation as pdf under the given filenmae. Needs PowerPoint installed.
        """
        return utils.save_as_pdf(self.prs, filename, overwrite)

    def save_as_png(self, save_folder, overwrite_folder=False) -> bool:
        """
        Saves the presentation as PNG's in the given folder. Needs PowerPoint installed.
        """
        return utils.save_as_png(self.prs, save_folder, overwrite_folder)
