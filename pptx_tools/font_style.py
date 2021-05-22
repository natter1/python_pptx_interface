"""
This module provides a helper class to deal with fonts in python-pptx.
@author: Nathanael Jöhrmann
"""
from typing import Union, Optional, Tuple

from pptx.dml.color import RGBColor
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE
from pptx.shapes.autoshape import Shape
from pptx.text.text import Font
from pptx.text.text import _Paragraph
from pptx.text.text import _Run
from pptx.util import Pt

from pptx_tools.enumerations import TEXT_CAPS_VALUES, TEXT_STRIKE_VALUES
from pptx_tools.fill_style import PPTXFillStyle
from pptx_tools.utils import _USE_DEFAULT, _DO_NOT_CHANGE


class PPTXFontStyle:
    """
    Helper class to deal with fonts in python-pptx. The internal class pptx.text.text.Font is limited, as it
    always needs an existing Text/Character/... for initializing and also basic functionality like assignment
    of one paragraph to another is missing.
    """
    # default language and paragraph; no _USE_DEFAULT for language_id -> use MSO_LANGUAGE_ID.NONE
    language_id: Union[MSO_LANGUAGE_ID, _USE_DEFAULT, None] = MSO_LANGUAGE_ID.ENGLISH_UK  # MSO_LANGUAGE_ID.GERMAN
    name: Union[str, _USE_DEFAULT, None] = "Roboto"  # "Arial"  # "Arial Narrow"

    def __init__(self):
        #  If set to use_default(), the bold, italic ... setting is cleared and is inherited
        #  from an enclosing shape’s setting, or a setting in a style or master
        self.bold: Union[bool, _USE_DEFAULT, None] = None
        self.italic: Union[bool, _USE_DEFAULT, None] = None
        self.underline: Union[MSO_TEXT_UNDERLINE_TYPE, _USE_DEFAULT, bool, None] = None

        # use class attribute; instance attribute only when changed by user
        # self.language_id: MSO_LANGUAGE_ID = MSO_LANGUAGE_ID.NONE  # ENGLISH_UK; ENGLISH_US; ESTONIAN; GERMAN; ...
        # self.name: Union[str, _USE_DEFAULT, None] = None

        # saved in units of Pt (not EMU like pptx.text.text.Font) - converting to EMU is done during write_to_font
        self.size: Optional[int] = None  # 18

        # todo: color is ColorFormat object
        self._color_rgb: Optional[RGBColor] = None
        # fil.fore_color changes paragraph color; also gradient or image might be useful (not implemented in FillStyle jet)
        self.fill_style: Optional[PPTXFillStyle] = None  # PPTXFillStyle()

        # experimental (not implemented in python-pptx):
        self.caps: Optional[TEXT_CAPS_VALUES] = None
        self.strikethrough: Optional[TEXT_STRIKE_VALUES] = None

    @property
    def color_rgb(self):
        return self._color_rgb

    @color_rgb.setter
    def color_rgb(self, value: Union[RGBColor, Tuple[any, any, any], None]):
        assert isinstance(value, RGBColor) or isinstance(value, tuple) or (value is None)
        self._color_rgb = RGBColor(*value) if isinstance(value, tuple) else value

    def read_font(self, font: Font) -> 'PPTXFontStyle':  # todo: check for None behavior (use_dfault() ? )
        """Read attributes from a pptx.text.text.Font object."""
        self.bold = font.bold
        self.italic = font.italic
        self.name = font.name
        self.size = None if font.size is None else font.size.pt
        self.underline = font.underline
        try:
            self.caps = TEXT_CAPS_VALUES(font._element.attrib['cap'])
        except KeyError:
            self.caps = None
        try:
            self.strikethrough = TEXT_STRIKE_VALUES(font._element.attrib['strike'])
        except KeyError:
            self.strikethrough = None
        return self

    def write_font(self, font: Font) -> None:
        """Write attributes to a pptx.text.text.Font object."""
        font.name = self._get_write_value(new_value=self.name, old_value=font.name)
        font.bold = self._get_write_value(new_value=self.bold, old_value=font.bold)
        font.italic = self._get_write_value(new_value=self.italic, old_value=font.italic)
        font.underline = self._get_write_value(new_value=self.underline, old_value=font.underline)

        if self.language_id == _USE_DEFAULT:
            font.language_id = MSO_LANGUAGE_ID.NONE
        else:
            font.language_id = self._get_write_value(new_value=self.language_id, old_value=font.language_id)

        if self.size is not None:
            font.size = None if self.size == _USE_DEFAULT else Pt(self.size)
        if self.color_rgb is not None:
            font.color.rgb = self.color_rgb

        if self.fill_style is not None:
            self.fill_style.write_fill(font.fill)

        self._write_caps(font)
        self._write_strikethrough(font)

    def _write_caps(self, font: Font):
        if self.caps is None:
            return
        else:
            font._element.attrib['cap'] = self.caps.value

    def _write_strikethrough(self, font: Font):
        if self.strikethrough is None:
            return
        else:
            font._element.attrib['strike'] = self.strikethrough.value

    @staticmethod
    def _get_write_value(new_value, old_value, check_default=True):
        """Used to check for None and use_default(), returning the correct value to write."""
        if new_value is None:
            return old_value
        if check_default and (new_value == _USE_DEFAULT):
            return None
        return new_value

    def write_shape(self, shape: Shape) -> None:
        """
        Write attributes to all paragraphs in given pptx.shapes.autoshape.Shape.
        Raises TypeError if given shape has no text_frame or table.
        """
        if shape.has_text_frame:
            self.write_text_frame(shape.text_frame)
        elif shape.has_table:
            for cell in shape.table.iter_cells():
                if not cell.is_spanned:
                    self.write_text_frame(cell.text_frame)
        else:
            raise TypeError("Cannot write paragraph for given shape (has no text_frame or table)")

    def write_text_frame(self, text_frame):
        """
        Write attributes to all paragraphs in given text_frame.
        """
        for paragraph in text_frame.paragraphs:
            self.write_paragraph(paragraph)

    def write_paragraph(self, paragraph: _Paragraph) -> None:
        """ Write attributes to given paragraph"""
        self.write_font(paragraph.font)

    def write_run(self, run: _Run) -> None:
        """ Write attributes to given run"""
        self.write_font(run.font)

    def set(self, bold: Optional[bool] = _DO_NOT_CHANGE,
            italic: Optional[bool] = _DO_NOT_CHANGE,
            language_id: MSO_LANGUAGE_ID = _DO_NOT_CHANGE,
            name: Optional[str] = _DO_NOT_CHANGE,
            size: Optional[int] = _DO_NOT_CHANGE,
            underline: Union[MSO_TEXT_UNDERLINE_TYPE, bool, None] = _DO_NOT_CHANGE,
            color_rgb: Union[RGBColor, Tuple[any, any, any]] = _DO_NOT_CHANGE,
            caps: Optional[TEXT_CAPS_VALUES] = _DO_NOT_CHANGE,
            strikethrough: Optional[TEXT_STRIKE_VALUES] = _DO_NOT_CHANGE
            ) -> 'PPTXFontStyle':
        """Convenience method to set several paragraph attributes together."""
        if bold is not _DO_NOT_CHANGE:
            self.bold = bold
        if italic is not _DO_NOT_CHANGE:
            self.italic = italic
        if language_id is not _DO_NOT_CHANGE:
            self.language_id = language_id
        if name is not _DO_NOT_CHANGE:
            self.name = name
        if size is not _DO_NOT_CHANGE:
            self.size = size
        if underline is not _DO_NOT_CHANGE:
            self.underline = underline
        if color_rgb is not _DO_NOT_CHANGE:
            self.color_rgb = color_rgb
        if caps is not _DO_NOT_CHANGE:
            self.caps = caps
        if strikethrough is not _DO_NOT_CHANGE:
            self.strikethrough = strikethrough
        return self

        # -----------------------------------------------------------------------------------------------
        # ----------------------------------- experimentell methods -------------------------------------
        # -----------------------------------------------------------------------------------------------
    # def _write_font_experimentell(self, paragraph: Font,all_caps: bool = True, strikethrough: bool = True):
    #     if all_caps:
    #         paragraph._element.attrib['cap'] = "all"
    #     else:
    #         pass
    #
    #     if strikethrough:
    #         paragraph._element.attrib['strike'] = "sngStrike"
    #     else:
    #         pass
