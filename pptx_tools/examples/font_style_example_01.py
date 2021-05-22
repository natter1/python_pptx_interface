"""
This script demonstrates how to work with fonts and paragraph-styles in python-pptx-interface.
@author: Nathanael Jöhrmann
"""

import os

from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE

from pptx_tools.creator import PPTXCreator
from pptx_tools.enumerations import TEXT_STRIKE_VALUES, TEXT_CAPS_VALUES
from pptx_tools.fill_style import PPTXFillStyle, FillType
from pptx_tools.font_style import PPTXFontStyle
from pptx_tools.position import PPTXPosition
from pptx_tools.style_sheets import font_title
from pptx_tools.templates import TemplateExample


def run(save_dir: str):
    filename_pptx = os.path.join(save_dir, "font_style_example_01.pptx")
    pp = PPTXCreator(TemplateExample())

    # default language and paragraph-type for all created PPTXFontStyle instances:
    PPTXFontStyle.lanaguage_id = MSO_LANGUAGE_ID.ENGLISH_UK
    PPTXFontStyle.name = "Roboto"

    title_slide = pp.add_title_slide("Font style example presentation")
    font = font_title()  # returns a PPTXFontStyle instance with bold paragraph and size = 32 Pt
    font.write_shape(title_slide.shapes.title)  # change paragraph attributes for all paragraphs in shape

    text_01 = "This text has four paragraphs. This is the first.\n" \
              "Das ist der zweite ...\n" \
              "... the third ...\n" \
              "... and the last."

    my_font = PPTXFontStyle()
    my_font.size = 16
    text_shape_01 = pp.add_text_box(title_slide, text_01, PPTXPosition(0.02, 0.24), my_font)

    my_font.set(size=22, bold=True, language_id=MSO_LANGUAGE_ID.GERMAN,
                strikethrough=TEXT_STRIKE_VALUES.SingleStrike,
                caps=TEXT_CAPS_VALUES.All)

    my_font.write_paragraph(text_shape_01.text_frame.paragraphs[1])

    my_font.set(size=18, bold=False, italic=True, name="Vivaldi",
                language_id=MSO_LANGUAGE_ID.ENGLISH_UK,
                underline=MSO_TEXT_UNDERLINE_TYPE.WAVY_DOUBLE_LINE,
                color_rgb=(255, 0, 0),
                strikethrough=None,
                caps=None)

    my_font.write_paragraph(text_shape_01.text_frame.paragraphs[2])

    my_font = PPTXFontStyle()
    my_font.set(size=52, bold=True)
    my_fill = PPTXFillStyle()
    my_fill.fill_type = FillType.PATTERNED
    my_fill.fore_color_rgb = (255, 0, 0)
    my_fill.back_color_rgb = (0, 0, 255)
    from pptx.enum.dml import MSO_PATTERN_TYPE
    my_fill.pattern = MSO_PATTERN_TYPE.PERCENT_50
    my_font.fill_style = my_fill

    my_font.write_paragraph(text_shape_01.text_frame.paragraphs[3])

    text_02 = "This text uses copied paragraph."

    my_copied_font = PPTXFontStyle()
    my_copied_font.read_font(text_shape_01.text_frame.paragraphs[1].font)
    text_shape_02 = pp.add_text_box(title_slide, text_02, PPTXPosition(0.42, 0.24), my_copied_font)

    pp.save(filename_pptx, overwrite=True)


if __name__ == '__main__':
    save_dir = os.path.dirname(os.path.abspath(__file__)) + '\\output\\'
    if not os.path.exists(save_dir):
        os.makedirs(save_dir)
    run(save_dir)
