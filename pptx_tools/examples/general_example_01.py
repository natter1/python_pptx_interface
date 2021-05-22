"""
This script demonstrates basic features of python-pptx-interface.
@author: Nathanael Jöhrmann
"""

import os

from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE, PP_PARAGRAPH_ALIGNMENT

from pptx_tools.creator import PPTXCreator
# from pptx_tools.fill_style import PPTXFillStyle, FillType
from pptx_tools.enumerations import TEXT_STRIKE_VALUES, TEXT_CAPS_VALUES
from pptx_tools.font_style import PPTXFontStyle
from pptx_tools.paragraph_style import PPTXParagraphStyle
from pptx_tools.position import PPTXPosition
from pptx_tools.style_sheets import font_title, font_default
from pptx_tools.templates import TemplateExample

try:
    import matplotlib.pyplot as plt

    matplotlib_installed = True
except ImportError as e:
    matplotlib_installed = False


def create_demo_figure():
    if not matplotlib_installed:
        return

    figure: plt.Figure = plt.figure(figsize=(3.4, 1.8), dpi=100, facecolor='w', edgecolor='w', frameon=True)
    figure.patch.set_alpha(0.5)
    sub_title = figure.suptitle('matplotlib figure', fontsize=14, fontweight='bold', color='red')
    sub_title.set_color('green')
    sub_title.set_rotation(5)
    sub_title.set_size(18)

    textstr = '\n'.join((
        fr'$\mu={5}^{5}$',
        r'$\mathrm{median}=3_3$',
        r'$median=3_3$',
        fr'$\sigma=$'))
    figure.text(0.3, 0.05, textstr)
    return figure


def run(save_dir: str):
    pp = PPTXCreator(TemplateExample())

    PPTXFontStyle.lanaguage_id = MSO_LANGUAGE_ID.ENGLISH_UK
    PPTXFontStyle.name = "Roboto"

    title_slide = pp.add_title_slide("General example 01 - title slide")
    font = font_title()  # returns a PPTXFontStyle instance with bold paragraph and size = 32 Pt
    font.write_shape(title_slide.shapes.title)  # change paragraph attributes for all paragraphs in shape

    slide2 = pp.add_slide("General example 01 - page2")
    pp.add_slide("General example 01 - page3")
    pp.add_slide("General example 01 - page4")
    pp.add_content_slide()  # add slide with hyperlinks to all other slides

    text = "This text has three paragraphs. This is the first.\n" \
           "Das ist der zweite ...\n" \
           "... and the third."
    my_font = font_default()
    my_font.size = 16
    text_shape = pp.add_text_box(title_slide, text, PPTXPosition(0.02, 0.24), my_font)

    my_font.set(size=22, bold=True, language_id=MSO_LANGUAGE_ID.GERMAN,
                strikethrough=TEXT_STRIKE_VALUES.SingleStrike,
                caps=TEXT_CAPS_VALUES.All)

    my_font.write_paragraph(text_shape.text_frame.paragraphs[1])

    my_font.set(size=18, bold=False, italic=True, name="Vivaldi",
                language_id=MSO_LANGUAGE_ID.ENGLISH_UK,
                underline=MSO_TEXT_UNDERLINE_TYPE.WAVY_DOUBLE_LINE,
                color_rgb=(255, 0, 0),
                strikethrough=None,
                caps=None)

    my_font.write_paragraph(text_shape.text_frame.paragraphs[2])

    table_data = []
    table_data.append([1, 2])  # rows can have different length
    table_data.append([4, slide2, 6])  # there is specific type needed for entries (implemented as text=f"{entry}")
    table_data.append(["", 8, 9])

    table = pp.add_table(title_slide, table_data, PPTXPosition(0.02, 0.4))
    paragraph_style = PPTXParagraphStyle()
    paragraph_style.set(alignment=PP_PARAGRAPH_ALIGNMENT.CENTER)
    paragraph_style.write_shape(table)

    if matplotlib_installed:
        fig = create_demo_figure()
        pp.add_matplotlib_figure(fig, title_slide, PPTXPosition(0.3, 0.4))
        pp.add_matplotlib_figure(fig, title_slide, PPTXPosition(0.3, 0.4, fig.get_figwidth(), -1.0), zoom=0.4)
        pp.add_matplotlib_figure(fig, title_slide, PPTXPosition(0.3, 0.4, fig.get_figwidth(), 0.0), zoom=0.5)
        pp.add_matplotlib_figure(fig, title_slide, PPTXPosition(0.3, 0.4, fig.get_figwidth(), 1.5), zoom=0.6)

        pp.add_text_box(title_slide, "Use latex-like syntax \nto create formula:", PPTXPosition(0.748, 0.23))
        pp.add_latex_formula(f"\mu={5}^{5}", title_slide, PPTXPosition(0.75, 0.35))
        formula02 = "\\int_0^\\infty e^{-x^2} dx=\\frac{\\sqrt{\\pi}}{2}"
        pp.add_latex_formula(formula02, title_slide, PPTXPosition(0.75, 0.45))
        pp.add_latex_formula(formula02, title_slide, PPTXPosition(0.75, 0.55), font_size=24, color="red")
        formula03 = "\\hat{x}, \\check{x}, \\tilde{a}, \\bar{\\ell}, \\dot{y}, \\ddot{y}, \\vec{z_1}, \\vec{z}_1"
        pp.add_latex_formula(formula03, title_slide, PPTXPosition(0.75, 0.65), font_size=24, color="blue")
        formula04 = r"\frac{3}{4} \binom{3}{4} \genfrac{}{}{0}{}{3}{4}"
        pp.add_latex_formula(formula04, title_slide, PPTXPosition(0.75, 0.75), font_size=44, color="g")
    pp.save(os.path.join(save_dir, "general_example_01.pptx"))

    try:  # only on Windows with PowerPoint installed:
        filename_pptx = os.path.join(save_dir, "general_example_01.pptx")
        filename_pdf = os.path.join(save_dir, "general_example_01.pdf")
        foldername_png = os.path.join(save_dir, "general_example_01_pngs")

        # use absolute path, because its not clear where PowerPoint saves PDF/PNG ... otherwise
        pp.save(filename_pptx, create_pdf=True, overwrite=True)
        pp.save_as_pdf(filename_pdf, overwrite=True)
        pp.save_as_png(foldername_png, overwrite_folder=True)
    except Exception as e:
        print(e)


if __name__ == '__main__':
    save_dir = os.path.dirname(os.path.abspath(__file__)) + '\\output\\'
    if not os.path.exists(save_dir):
        os.makedirs(save_dir)
    run(save_dir)
