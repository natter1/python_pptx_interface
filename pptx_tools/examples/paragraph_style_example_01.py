"""
This script demonstrates how to work with paragraph-styles in python-pptx-interface.
@author: Nathanael Jöhrmann
"""
import os

from pptx.util import Inches

from pptx_tools.creator import PPTXCreator
from pptx_tools.paragraph_style import PPTXParagraphStyle
from pptx_tools.position import PPTXPosition
from pptx_tools.templates import TemplateExample
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def run(save_dir: str):
    pp = PPTXCreator(TemplateExample())
    title_slide = pp.add_title_slide("Paragraph style example 01 - title slide")

    text = "This text has three paragraphs. This is the first.\n" \
           "Das ist der zweite ...\n" \
           "... and the third."
    text_shape = pp.add_text_box(title_slide, text, PPTXPosition(0.02, 0.24))
    text_shape.width = Inches(6)
    text_shape.height = Inches(2)
    paragraph_style = PPTXParagraphStyle()
    paragraph_style.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph_style.write_paragraph(text_shape.text_frame.paragraphs[0])
    paragraph_style.line_spacing = 3
    paragraph_style.write_paragraph(text_shape.text_frame.paragraphs[1])
    paragraph_style.level = 8
    paragraph_style.write_paragraph(text_shape.text_frame.paragraphs[2])

    filename_pptx = os.path.join(save_dir, "paragraph_style_example_01.pptx")
    pp.save(filename_pptx, create_pdf=True, overwrite=True)
    # foldername_png = os.path.join(save_dir, "paragraph_style_example_01_pngs")
    # pp.save_as_png(foldername_png, overwrite_folder=True)


if __name__ == '__main__':
    save_dir = os.path.dirname(os.path.abspath(__file__)) + '\\output\\'
    if not os.path.exists(save_dir):
        os.makedirs(save_dir)
    run(save_dir)