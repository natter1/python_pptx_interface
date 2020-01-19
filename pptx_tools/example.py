import os

from pptx_tools.creator import PPTXCreator, PPTXPosition
from pptx_tools.style_sheets import font_title, font_default
from pptx_tools.templates import TemplateExample

from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE

try:
    import matplotlib.pyplot as plt
    matplotlib_installed = True
except ImportError as e:
    matplotlib_installed = False


def create_demo_figure():
    if not matplotlib_installed:
        return

    figure: plt.Figure = plt.figure(figsize=(3.4, 1.8), dpi=100, facecolor='w', edgecolor='w', frameon=True)
    figure.patch.set_alpha(0.5)
    sub_title = figure.suptitle('matplotlib figure', fontsize=14, fontweight='bold', color='red')
    sub_title.set_color('green')
    sub_title.set_rotation(5)
    sub_title.set_size(18)

    textstr = '\n'.join((
        fr'$\mu={5}^{5}$',
        r'$\mathrm{median}=3_3$',
        r'$median=3_3$',
        fr'$\sigma=$'))
    figure.text(0.3, 0.05, textstr)
    return figure


def run(save_dir: str):
    pp = PPTXCreator(TemplateExample())

    title_slide = pp.add_title_slide("Example presentation")
    font = font_title()
    font.write_shape(title_slide.shapes.title)  # change font attributes for all paragraphs in shape

    slide2 = pp.add_slide("page2")
    pp.add_slide("page3")
    pp.add_slide("page4")
    pp.add_content_slide()

    text = "This text has three paragraphs. This is the first.\n" \
           "Das ist der zweite ...\n" \
           "... and the third."
    my_font = font_default()
    my_font.size = 16
    text_shape = pp.add_text_box(title_slide, text, PPTXPosition(0.02, 0.24), my_font)

    my_font.set(size=22, bold=True, language_id=MSO_LANGUAGE_ID.GERMAN)
    my_font.write_paragraph(text_shape.text_frame.paragraphs[1])

    my_font.set(size=18, bold=False, italic=True, name="Vivaldi",
                language_id=MSO_LANGUAGE_ID.ENGLISH_UK,
                underline=MSO_TEXT_UNDERLINE_TYPE.WAVY_DOUBLE_LINE)
    my_font.write_paragraph(text_shape.text_frame.paragraphs[2])

    table_data = []
    table_data.append([1, 2])  # rows can have different length
    table_data.append([4, slide2, 6])  # there is specific type needed for entries (implemented as text=f"{entry}")
    table_data.append(["", 8, 9])

    pp.add_table(title_slide, table_data, PPTXPosition(0.02, 0.4))


    if matplotlib_installed:
        fig = create_demo_figure()
        pp.add_matplotlib_figure(fig, title_slide, PPTXPosition(0.3, 0.4))
        pp.add_matplotlib_figure(fig, title_slide, PPTXPosition(0.3, 0.4, fig.get_figwidth(), -1.0), zoom=0.4)
        pp.add_matplotlib_figure(fig, title_slide, PPTXPosition(0.3, 0.4, fig.get_figwidth(), 0.0), zoom=0.5)
        pp.add_matplotlib_figure(fig, title_slide, PPTXPosition(0.3, 0.4, fig.get_figwidth(), 1.5), zoom=0.6)

    pp.save(os.path.join(save_dir, "example.pptx"))

    try:  # only on Windows with PowerPoint installed:
        filename_pptx = os.path.join(save_dir, "example.pptx")
        filename_pdf = os.path.join(save_dir, "example.pdf")
        foldername_png = os.path.join(save_dir, "example_pngs")

        # use absolute path, because its not clear where PowerPoint saves PDF/PNG ... otherwise
        pp.save(filename_pptx, create_pdf=True, overwrite=True)
        pp.save_as_pdf(filename_pdf, overwrite=True)
        pp.save_as_png(foldername_png, overwrite_folder=True)
    except Exception as e:
        print(e)


if __name__ == '__main__':
    run(os.path.dirname(os.path.abspath(__file__))
)
