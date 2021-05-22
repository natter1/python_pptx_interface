"""
This file contains variables with names of important pptx template master_slide shapes
"""
from datetime import datetime

import pkg_resources
# from pptx.enum.text import MSO_AUTO_SIZE
from pptx import Presentation

from pptx_tools.better_abc import ABCMeta, abstract_attribute
from pptx_tools.utils import change_paragraph_text_to


class AbstractTemplate(metaclass=ABCMeta):
    """
    Templates should subclass this abstract class, to make sure,
    that all important attributes/methodes are defined.
    """

    @abstract_attribute
    def TEMPLATE_FILE(cls):
        pass

    @abstract_attribute
    def prs(self):
        pass

    @abstract_attribute
    def title_layout(cls):
        pass

    @abstract_attribute
    def default_layout(cls):
        pass


# ----------------------------------------------------------------------------
# --------- Customized template classes are needed for each template ---------
# ----------------------------------------------------------------------------
class TemplateExample(AbstractTemplate):
    """
    Class handling example-template.pptx).
    """
    TEMPLATE_FILE = pkg_resources.resource_filename('pptx_tools', 'resources/example-template.pptx')

    def __init__(self):
        self.prs = Presentation(self.TEMPLATE_FILE)

        self.title_layout = self.prs.slide_masters[0].slide_layouts[0]
        self.default_layout = self.prs.slide_masters[1].slide_layouts[0]

        # following names are the same for both slide masters
        self.author_shape_name = "Rectangle 4"
        self.website_shape_name = "Rectangle 5"

        date_time = datetime.now().strftime("%d %B, %Y")
        self.set_author("Nathanael Jöhrmann", city="Chemnitz", date=date_time)
        self.set_website("https://github.com/natter1/python_pptx_interface")

    def set_author(self, name, city=None, date=None):
        text = ""
        spacer = " ∙ "
        if city:
            text += city + spacer
        if date:
            text += date + spacer
        text += name
        self.write_text_to_master_shape(text=text, shape_name=self.author_shape_name)

    def set_website(self, text):
        self.write_text_to_master_shape(text=text, shape_name=self.website_shape_name)

    def write_text_to_master_shape(self, text, shape_name):
        for shape in self.master_shapes:
            if not shape.has_text_frame:
                continue
            if shape.name == shape_name:
                change_paragraph_text_to(shape.text_frame.paragraphs[0], text)

    @property
    def master_shapes(self):
        result = []
        for slide_master in self.prs.slide_masters:
            result.extend(slide_master.shapes)
        return result


def analyze_pptx(template_file):
    """ Take the given file and analyze the structure of master slides.
    Prints shape names/ids and texts for SlideMaster-shapes
    To get an output file contains marked up information
    remove comment on last two lines of function.
    This is helpful when manipulating template-files.
    """
    prs = Presentation(template_file)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    slide_masters = prs.slide_masters
    for index, slide_master in enumerate(prs.slide_masters):
        print('------------------------------------')
        print('------------------------------------')
        print(f"slide master indexed: {index}")
        print(slide_master)
        print("text boxes:")
        for shape in slide_master.shapes:
            try:
                dummystring = f"shape name: {shape.name} - shape text: {shape.text}"
                shape.text = shape.name
                print(dummystring)
            except:
                pass
            # shape.text = 'hahahaha'
        # for shape in slide_master.slideshapes:
        #     print(shape)
        print('------------------------------------')
        for index, slide_layout in enumerate(slide_master.slide_layouts):
            print(f"\tslide layout: {slide_layout.name}")
            slide = prs.slides.add_slide(slide_master.slide_layouts[index])
            # Not every slide has to have a title
            try:
                title = slide.shapes.title
                title.text = 'Title for Layout {}'.format(index)
            except AttributeError:
                print("No Title for Layout {}".format(index))
            # Go through all the placeholders and identify them by index and type
            for shape in slide.placeholders:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    # Do not overwrite the title which is just a special placeholder
                    try:
                        if 'Title' not in shape.text:
                            shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                    except AttributeError:
                        print("{} has no text attribute".format(phf.type))
                    print(f"\t\tid: {phf.idx} - name: {shape.name}")
    # output_file = '..\\resources\pptx\\template_names.pptx'
    # prs.save(output_file)


def analyze_paragraphs(paragraphs):
    for index, para in enumerate(paragraphs):
        print(f"index: {index} - text: {para.text}")
        for run_index, run in enumerate(para.runs):
            print(f"\trun: {run.text}")
