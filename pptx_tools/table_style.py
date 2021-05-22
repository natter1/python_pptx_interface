"""
This module provides a helper class to deal with tables in python-pptx.
@author: Nathanael Jöhrmann
"""
from typing import Optional

from pptx.shapes.autoshape import Shape
from pptx.table import Table, _Cell
from pptx.util import Inches

from pptx_tools.fill_style import PPTXFillStyle
from pptx_tools.font_style import PPTXFontStyle
from pptx_tools.position import PPTXPosition
from pptx_tools.utils import iter_table_cells, _DO_NOT_CHANGE


class PPTXCellStyle:  # format table cell
    def __init__(self):
        self.fill_style = PPTXFillStyle()

    def write_cell(self, cell: _Cell) -> None:
        self.fill_style.write_fill(cell.fill)


class PPTXTableStyle:
    """
    ...
    """

    def __init__(self):
        self.font_style: Optional[PPTXFontStyle] = None  # PPTXFontStyle()
        self.cell_style: Optional[PPTXCellStyle] = None  # PPTXCellStyle()
        self.first_row_header: Optional[bool] = None  # False  # special formatting for first row?
        self.col_banding: Optional[bool] = None  # False  # slightly alternate color brightness per col
        self.row_banding: Optional[bool] = None  # True  # slightly alternate color brightness per row

        self.width: Optional[float] = None  # in [Inches]; don't use Inches() - is transformed in _write_col_sizes!!!
        self.col_ratios = None
        self.position = None

    def read_table(self, table: Table):
        """Read attributes from a Table object, ignoring font and cell style."""
        self.first_row_header = table._tbl.firstRow
        self.col_banding = table.vert_banding
        self.row_banding = table.horz_banding

        self.col_ratios = []
        for column in table.columns:
            self.col_ratios.append(column.width.inches)
        self.width = sum(self.col_ratios)

    def set(self, font_style: Optional[PPTXFontStyle] = _DO_NOT_CHANGE,
            cell_style: Optional[PPTXCellStyle] = _DO_NOT_CHANGE,
            first_row_header: Optional[bool] = _DO_NOT_CHANGE,
            col_banding: Optional[bool] = _DO_NOT_CHANGE,
            row_banding: Optional[bool] = _DO_NOT_CHANGE,
            width: Optional[float] = _DO_NOT_CHANGE,
            col_ratios: Optional[list] = _DO_NOT_CHANGE,
            position: Optional[PPTXPosition] = _DO_NOT_CHANGE
            ) -> 'PPTXTableStyle':
        """Convenience method to set several table attributes together."""
        if font_style is not _DO_NOT_CHANGE:
            self.font_style = font_style
        if cell_style is not _DO_NOT_CHANGE:
            self.cell_style = cell_style
        if first_row_header is not _DO_NOT_CHANGE:
            self.first_row_header = first_row_header
        if col_banding is not _DO_NOT_CHANGE:
            self.col_banding = col_banding
        if row_banding is not _DO_NOT_CHANGE:
            self.row_banding = row_banding
        if width is not _DO_NOT_CHANGE:
            self.width = width
        if col_ratios is not _DO_NOT_CHANGE:
            self.col_ratios = col_ratios
        if position is not _DO_NOT_CHANGE:
            self.position = position
        return self

    def _write_all_cells(self, table: Table) -> None:
        for cell in iter_table_cells(table):
            if self.font_style is not None:
                # paragraph is managed per cell; there is no "table paragraph"
                self.font_style.write_text_frame(cell.text_frame)
            if self.cell_style is not None:
                self.cell_style.write_cell(cell)

    def _update_col_ratios(self, number_of_cols: int) -> None:
        """Add default values (1) if col_ratios has not enough entries for all table cols."""
        if self.col_ratios is None:
            self.col_ratios = []
        while len(self.col_ratios) < number_of_cols:
            self.col_ratios.append(1)

    def _write_col_sizes(self, table):
        assert self.width is not None
        number_of_cols = len(table.columns)
        self._update_col_ratios(number_of_cols)

        ratio_sum = sum(self.col_ratios[:len(table.columns)])
        for column, ratio in zip(table.columns, self.col_ratios):
            column.width = Inches(self.width * ratio / ratio_sum)

    def write_shape(self, shape: Shape) -> None:
        """Write attributes to table in given pptx.shapes.autoshape.Shape."""
        if not shape.has_table:
            print(f"Warning: Could not write table style. {shape} has no table.")
            return
        if self.position is not None:
            shape.left, shape.top = self.position.tuple()
        self.write_table(shape.table)

    def write_table(self, table: Table) -> None:
        """Write attributes to table object."""
        if self.first_row_header is not None:
            table._tbl.firstRow = self.first_row_header

        if self.col_banding is not None:
            table.vert_banding = self.col_banding

        if self.row_banding is not None:
            table.horz_banding = self.row_banding

        if self.width is not None:
            self._write_col_sizes(table)
        self._write_all_cells(table)

    def set_width_as_fraction(self, fraction: float):
        assert fraction > 0.0
        if PPTXPosition.prs is None:
            raise TypeError("Still no presentation set for PPTXPosition."
                            " Create a PPTXCreator instance first, or set manually.")

        self.width = PPTXPosition.prs.slide_width.inches * fraction
