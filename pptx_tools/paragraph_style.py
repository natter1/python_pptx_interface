"""
This module provides a helper class to deal with paragraphs in python-pptx.
@author: Nathanael Jöhrmann
"""
from typing import Optional

from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.shapes.autoshape import Shape
from pptx.text.text import _Paragraph

from pptx_tools.font_style import PPTXFontStyle
from pptx_tools.utils import _DO_NOT_CHANGE


class PPTXParagraphStyle:
    """
    Helper class to deal with paragraphs in python-pptx.
    """

    def __init__(self):
        # the following values cannot be inherited from parent object(?) -----------------------------------------------
        self.alignment: Optional[PP_PARAGRAPH_ALIGNMENT] = None  # PP_PARAGRAPH_ALIGNMENT.CENTER/JUSTIFY/LEFT/RIGHT/...
        self.level: Optional[int] = None  # 0 .. 8 (indentation level)
        self.line_spacing: Optional[float] = None
        self.space_before: Optional[float] = None
        self.space_after: Optional[float] = None
        # --------------------------------------------------------------------------------------------------------------
        self.font_style: Optional[PPTXFontStyle] = None

    def read_paragraph(self, paragraph: _Paragraph) -> 'PPTXParagraphStyle':
        """Read attributes from a _Paragraph object."""
        self.alignment = paragraph.alignment
        self.level = paragraph.level
        self.line_spacing = paragraph.line_spacing
        self.space_before = paragraph.space_before
        self.space_after = paragraph.space_after
        self.font_style = PPTXFontStyle().read_font(paragraph.font)
        return self

    def set(self, alignment: Optional[PP_PARAGRAPH_ALIGNMENT] = _DO_NOT_CHANGE,
            level: Optional[int] = _DO_NOT_CHANGE,
            line_spacing: Optional[float] = _DO_NOT_CHANGE,
            space_before: Optional[float] = _DO_NOT_CHANGE,
            space_after: Optional[float] = _DO_NOT_CHANGE
            ) -> 'PPTXParagraphStyle':
        """Convenience method to set several paragraph attributes together."""
        if alignment is not _DO_NOT_CHANGE:
            self.alignment = alignment
        if level is not _DO_NOT_CHANGE:
            self.level = level
        if line_spacing is not _DO_NOT_CHANGE:
            self.line_spacing = line_spacing
        if space_before is not _DO_NOT_CHANGE:
            self.space_before = space_before
        if space_after is not _DO_NOT_CHANGE:
            self.space_after = space_after
        return self

    def write_paragraph(self, paragraph: _Paragraph) -> None:
        """Write paragraph style to given paragraph."""
        if self.alignment is not None:
            paragraph.alignment = self.alignment
        if self.level is not None:
            paragraph.level = self.level
        if self.line_spacing is not None:
            paragraph.line_spacing = self.line_spacing
        if self.font_style is not None:
            self.font_style.write_paragraph(paragraph)

    def write_shape(self, shape: Shape) -> None:
        """
        Write attributes to all paragraphs in given pptx.shapes.autoshape.Shape.
        Raises TypeError if given shape has no text_frame or table.
        """
        if shape.has_text_frame:
            self.write_text_frame(shape.text_frame)
        elif shape.has_table:
            for cell in shape.table.iter_cells():
                if not cell.is_spanned:
                    self.write_text_frame(cell.text_frame)
        else:
            raise TypeError("Cannot write paragraph for given shape (has no text_frame or table)")

    def write_text_frame(self, text_frame):
        """
        Write attributes to all paragraphs in given text_frame.
        """
        for paragraph in text_frame.paragraphs:
            self.write_paragraph(paragraph)
