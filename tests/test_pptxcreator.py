"""
This file contains tests for all PPTXCreator-methods. It can be useful to open the created files in the local temp
folder, to also check functionality inside PowerPoint.
@author: Nathanael Jöhrmann
"""

import glob
import os

import matplotlib.pyplot as plt
import pytest
from pptx.util import Inches

from pptx_tools.creator import PPTXCreator
from pptx_tools.position import PPTXPosition
from pptx_tools.style_sheets import table_no_header
from pptx_tools.templates import TemplateExample


@pytest.fixture(scope='class')
def pptx_creator():
    creator = PPTXCreator(TemplateExample())
    yield creator


@pytest.fixture(scope="function", params=[-4, 0, 0.4, 1, 2])
def fraction_test_case(request):
    yield request.param


@pytest.fixture(scope='class')
def matplotlib_figure():
    yield plt.figure(figsize=(2, 1), dpi=100, facecolor='w', edgecolor='w', frameon=True)


class TestPPTXCreator:
    def test__create_presentation(self, pptx_creator):
        old_prs = pptx_creator.prs
        pptx_creator._create_presentation()
        new_prs = pptx_creator.prs
        assert old_prs is not new_prs
        assert new_prs is not None
        pptx_creator._create_presentation(TemplateExample())
        new_prs_with_template = pptx_creator.prs
        assert new_prs is not new_prs_with_template
        assert new_prs_with_template is not None

    def test__create_presentation_from_template(self, pptx_creator):
        old_prs = pptx_creator.prs
        pptx_creator._create_presentation_from_template(TemplateExample())
        new_prs_with_template = pptx_creator.prs
        assert old_prs is not new_prs_with_template
        assert new_prs_with_template is not None

    def test__fraction_height_to_inch(self, pptx_creator, fraction_test_case):
        inches_calc = pptx_creator._fraction_height_to_inch(fraction_test_case)
        inches_from_prs = pptx_creator.prs.slide_height * fraction_test_case
        assert type(inches_calc) is Inches  # type is Inches even if data is int!
        assert inches_calc == inches_from_prs

    def test__fraction_width_to_inch(self, pptx_creator, fraction_test_case):
        inches_calc = pptx_creator._fraction_width_to_inch(fraction_test_case)
        inches_from_prs = pptx_creator.prs.slide_width * fraction_test_case
        assert type(inches_calc) is Inches  # type is Inches even if data is int!
        assert inches_calc == inches_from_prs

    def test__get_rows_cols(self, pptx_creator):
        table_data = [[0, 1, 2], [1], [2], [3], [4]]  # 5 rows; 3 cols
        result = pptx_creator._get_rows_cols(table_data)
        assert result == (5, 3)

    def test_add_content_slide(self, pptx_creator):  # todo: how to improve test?
        slide = pptx_creator.add_content_slide()
        assert slide

    def test_add_latex_formula(self, pptx_creator):
        slide = pptx_creator.add_slide("test_add_latex_formula")
        pptx_creator.add_latex_formula("a=b", slide, PPTXPosition(0.25, 0.25))
        assert True  # todo: How to test correct formula creation in code?

    def test_add_matplotlib_figure(self, pptx_creator, matplotlib_figure):
        slide = pptx_creator.add_slide("test_add_matplotlib_figure")
        fig_width = matplotlib_figure.bbox_inches.width
        fig_height = matplotlib_figure.bbox_inches.height
        zoom = 0.8
        position = PPTXPosition(0.6, 0.4, 1, -1)
        shape = pptx_creator.add_matplotlib_figure(matplotlib_figure, slide, position, zoom=zoom)
        assert pptx_creator.prs.slide_width * position.left_rel + Inches(position.left) == shape.left
        assert pptx_creator.prs.slide_height * position.top_rel + Inches(position.top) == shape.top
        assert fig_width * zoom == shape.width.inches
        assert fig_height * zoom == shape.height.inches

    def test_add_slide(self, pptx_creator):
        n_slides_before = len(pptx_creator.prs.slides)
        pptx_creator.add_slide(" test_add_slide")
        n_slides_after = len(pptx_creator.prs.slides)
        assert n_slides_before == n_slides_after - 1

    def test_add_table(self, pptx_creator):
        table_data = [[0, 1, 2], [1], [2], [3], [4]]  # 5 rows; 3 cols
        position = PPTXPosition(0.6, 0.4, 1, -1)
        table_style = table_no_header()
        slide = pptx_creator.add_slide("test_add_table")
        shape = pptx_creator.add_table(slide, table_data=table_data, position=position, table_style=table_style)
        assert pptx_creator.prs.slide_width * position.left_rel + Inches(position.left) == shape.left
        assert pptx_creator.prs.slide_height * position.top_rel + Inches(position.top) == shape.top
        assert len(shape.table.rows) == 5

    def test_add_text_box(self, pptx_creator):
        slide = pptx_creator.add_slide("test_add_text_box")
        position = PPTXPosition(0.6, 0.4, 1, -1)
        text = "Test text"
        shape = pptx_creator.add_text_box(slide, text, position)
        assert pptx_creator.prs.slide_width * position.left_rel + Inches(position.left) == shape.left
        assert pptx_creator.prs.slide_height * position.top_rel + Inches(position.top) == shape.top
        assert text == shape.text

    def test_add_title_slide(self, pptx_creator):
        n_slides_before = len(pptx_creator.prs.slides)
        pptx_creator.add_title_slide("test_add_title_slide")
        n_slides_after = len(pptx_creator.prs.slides)
        assert n_slides_before == n_slides_after - 1

    def test_create_hyperlink(self, pptx_creator):  # todo: how to improve test?
        slide_from = pptx_creator.add_slide("test_create_hyperlink (from)")
        slide_to = pptx_creator.add_slide("test_create_hyperlink (to)")
        shape = pptx_creator.add_text_box(slide_from,
                                          "hyperlink to slide \"test_create_hyperlink (to)\"",
                                          PPTXPosition(0.25, 0.25))
        pptx_creator.create_hyperlink(shape.text_frame.paragraphs[0].runs[0], shape, slide_to)
        # pptx_creator.move_slide(slide_to, 0)
        assert True

    def test_move_slide(self, pptx_creator):
        slide_01 = pptx_creator.add_slide("test_move_slide_01 (not moved)")
        slide_02 = pptx_creator.add_slide("test_move_slide_02 (move to front)")
        pptx_creator.move_slide(slide_02, 0)
        assert pptx_creator.prs.slides._sldIdLst.sldId_lst[0].id != slide_01.slide_id
        assert pptx_creator.prs.slides._sldIdLst.sldId_lst[0].id == slide_02.slide_id

    def test_remove_unpopulated_shapes(self, pptx_creator):
        slide = pptx_creator.prs.slides.add_slide(pptx_creator.default_layout)
        assert len(slide.shapes) == 3  # make sure, default template wasn't changed
        slide.shapes.title.text = "test_remove_unpopulated_shapes"
        pptx_creator.remove_unpopulated_shapes(slide)
        assert len(slide.shapes) == 1  # (one shape populated: title)

    def test_save(self, pptx_creator, tmpdir):
        file = tmpdir.join("test_save.pptx")
        pptx_creator.save(file)
        assert os.path.isfile(file)
        old_size = os.path.getsize(file)
        pptx_creator.add_slide("test_save")  # increase file size
        pptx_creator.save(file)
        assert old_size == os.path.getsize(file)  # file not overwritten
        assert not os.path.isfile(tmpdir.join("test_save.pdf"))  # create_pdf should be False by default

    def test_save__create_pdf_is_true(self, pptx_creator, tmpdir):
        pptx_file = tmpdir.join("test_save.pptx")
        pdf_file = tmpdir.join("test_save.pdf")
        pptx_creator.add_slide("test_save__create_pdf_is_true")  # PowerPoint cannot export empty presentation as pdf
        pptx_creator.save(pptx_file, create_pdf=True)
        assert os.path.isfile(pptx_file)
        assert os.path.isfile(pdf_file)
        old_size_pdf = os.path.getsize(pdf_file)
        old_size_pptx = os.path.getsize(pptx_file)
        pptx_creator.add_slide("test_save__create_pdf_is_true_2 (should not be in saved files!)")  # increase file size
        pptx_creator.save(pptx_file, create_pdf=True)
        assert old_size_pdf == os.path.getsize(pdf_file)  # file not overwritten
        assert old_size_pptx == os.path.getsize(pptx_file)  # file not overwritten

    def test_save__overwrite_is_true(self, pptx_creator, tmpdir):
        pptx_file = tmpdir.join("test_save.pptx")
        pptx_creator.save(pptx_file, overwrite=True)
        assert os.path.isfile(pptx_file)
        old_size = os.path.getsize(pptx_file)
        pptx_creator.add_slide("test_save__overwrite_is_true")  # increase file size
        pptx_creator.save(pptx_file, overwrite=True)
        assert old_size < os.path.getsize(pptx_file)  # file overwritten

    def test_save_as_pdf(self, pptx_creator, tmpdir):
        pdf_file = tmpdir.join("test_save.pdf")
        pptx_creator.add_slide("test_save")  # PowerPoint cannot export empty presentation as pdf
        pptx_creator.save_as_pdf(pdf_file)
        assert os.path.isfile(pdf_file)
        old_size_pdf = os.path.getsize(pdf_file)
        pptx_creator.add_slide("test_save slide 2")  # increase file size
        pptx_creator.save_as_pdf(pdf_file)
        assert old_size_pdf == os.path.getsize(pdf_file)  # file not overwritten

    def test_save_as_pdf__overwrite_is_true(self, pptx_creator, tmpdir):
        pdf_file = tmpdir.join("test_save.pdf")
        pptx_creator.add_slide("test_save")  # PowerPoint cannot export empty presentation as pdf
        pptx_creator.save_as_pdf(pdf_file, overwrite=True)
        assert os.path.isfile(pdf_file)
        old_size_pdf = os.path.getsize(pdf_file)
        pptx_creator.add_slide("test_save slide 2")  # increase file size
        pptx_creator.save_as_pdf(pdf_file, overwrite=True)
        assert old_size_pdf < os.path.getsize(pdf_file)  # file not overwritten

    def test_save_as_png(self, pptx_creator, tmpdir):
        pptx_creator.add_slide("test_save")
        assert not pptx_creator.save_as_png(tmpdir)  # tmpdir exists -> no files saved
        assert len(glob.glob(f"{tmpdir}\\*.png")) == 0
        assert pptx_creator.save_as_png(tmpdir.join("pngs"))  # create new folder, or set overwrite = True
        assert len(glob.glob(f"{tmpdir.join('pngs')}\\*.png")) == len(pptx_creator.prs.slides)
